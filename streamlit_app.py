"""
AWS Enterprise Assessment Platform
Control Tower Migration & Golden Architecture (Serverless) Assessment Tool
AI-Driven Enterprise Grade Assessment Application
"""

import streamlit as st
import json
import os
from datetime import datetime
from typing import Dict, List, Any, Optional
import io
import base64

# Configure page
st.set_page_config(
    page_title="AWS Enterprise Assessment Platform",
    page_icon="🏗️",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for enterprise-grade styling
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=IBM+Plex+Sans:wght@300;400;500;600;700&family=IBM+Plex+Mono:wght@400;500&display=swap');
    
    :root {
        --primary-color: #0f1419;
        --secondary-color: #1a2634;
        --accent-color: #ff9500;
        --accent-secondary: #00d4aa;
        --text-primary: #ffffff;
        --text-secondary: #8899a6;
        --success-color: #00d4aa;
        --warning-color: #ff9500;
        --danger-color: #ff6b6b;
        --card-bg: #1a2634;
        --border-color: #2d3e50;
    }
    
    .stApp {
        background: linear-gradient(135deg, #0f1419 0%, #1a2634 50%, #0f1419 100%);
    }
    
    .main-header {
        background: linear-gradient(90deg, #1a2634 0%, #2d3e50 100%);
        padding: 2rem 2.5rem;
        border-radius: 16px;
        margin-bottom: 2rem;
        border: 1px solid #2d3e50;
        box-shadow: 0 8px 32px rgba(0,0,0,0.3);
    }
    
    .main-header h1 {
        font-family: 'IBM Plex Sans', sans-serif;
        font-weight: 700;
        font-size: 2.2rem;
        background: linear-gradient(90deg, #ff9500, #00d4aa);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        margin: 0;
        letter-spacing: -0.5px;
    }
    
    .main-header p {
        font-family: 'IBM Plex Sans', sans-serif;
        color: #8899a6;
        font-size: 1rem;
        margin-top: 0.5rem;
    }
    
    .assessment-card {
        background: linear-gradient(145deg, #1a2634 0%, #243447 100%);
        padding: 1.75rem;
        border-radius: 16px;
        border: 1px solid #2d3e50;
        margin-bottom: 1.5rem;
        transition: all 0.3s ease;
        box-shadow: 0 4px 20px rgba(0,0,0,0.2);
    }
    
    .assessment-card:hover {
        border-color: #ff9500;
        box-shadow: 0 8px 32px rgba(255,149,0,0.15);
        transform: translateY(-2px);
    }
    
    .metric-card {
        background: linear-gradient(145deg, #243447 0%, #1a2634 100%);
        padding: 1.5rem;
        border-radius: 12px;
        text-align: center;
        border: 1px solid #2d3e50;
    }
    
    .metric-value {
        font-family: 'IBM Plex Mono', monospace;
        font-size: 2.5rem;
        font-weight: 700;
        color: #00d4aa;
    }
    
    .metric-label {
        font-family: 'IBM Plex Sans', sans-serif;
        font-size: 0.85rem;
        color: #8899a6;
        text-transform: uppercase;
        letter-spacing: 1px;
        margin-top: 0.5rem;
    }
    
    .domain-header {
        font-family: 'IBM Plex Sans', sans-serif;
        font-weight: 600;
        font-size: 1.1rem;
        color: #ff9500;
        margin-bottom: 1rem;
        padding-bottom: 0.5rem;
        border-bottom: 2px solid #2d3e50;
    }
    
    .score-badge {
        display: inline-block;
        padding: 0.4rem 1rem;
        border-radius: 20px;
        font-family: 'IBM Plex Mono', monospace;
        font-weight: 600;
        font-size: 0.9rem;
    }
    
    .score-high { background: rgba(0,212,170,0.2); color: #00d4aa; border: 1px solid #00d4aa; }
    .score-medium { background: rgba(255,149,0,0.2); color: #ff9500; border: 1px solid #ff9500; }
    .score-low { background: rgba(255,107,107,0.2); color: #ff6b6b; border: 1px solid #ff6b6b; }
    
    .stButton > button {
        font-family: 'IBM Plex Sans', sans-serif;
        font-weight: 600;
        background: linear-gradient(90deg, #ff9500, #ff7b00);
        color: white;
        border: none;
        padding: 0.75rem 2rem;
        border-radius: 8px;
        transition: all 0.3s ease;
        box-shadow: 0 4px 15px rgba(255,149,0,0.3);
    }
    
    .stButton > button:hover {
        background: linear-gradient(90deg, #ff7b00, #ff5500);
        box-shadow: 0 6px 20px rgba(255,149,0,0.4);
        transform: translateY(-2px);
    }
    
    .stRadio > label, .stCheckbox > label {
        font-family: 'IBM Plex Sans', sans-serif;
        color: #ffffff !important;
    }
    
    .stSelectbox > label, .stTextArea > label, .stTextInput > label {
        font-family: 'IBM Plex Sans', sans-serif;
        color: #8899a6 !important;
        font-weight: 500;
    }
    
    .recommendation-card {
        background: linear-gradient(145deg, #1a2634 0%, #0f1419 100%);
        border-left: 4px solid #00d4aa;
        padding: 1.25rem;
        border-radius: 0 12px 12px 0;
        margin: 1rem 0;
    }
    
    .gap-card {
        background: linear-gradient(145deg, #1a2634 0%, #0f1419 100%);
        border-left: 4px solid #ff6b6b;
        padding: 1.25rem;
        border-radius: 0 12px 12px 0;
        margin: 1rem 0;
    }
    
    .insight-card {
        background: linear-gradient(145deg, #1a2634 0%, #0f1419 100%);
        border-left: 4px solid #ff9500;
        padding: 1.25rem;
        border-radius: 0 12px 12px 0;
        margin: 1rem 0;
    }
    
    .stTabs [data-baseweb="tab-list"] {
        gap: 8px;
        background-color: #1a2634;
        padding: 0.5rem;
        border-radius: 12px;
    }
    
    .stTabs [data-baseweb="tab"] {
        font-family: 'IBM Plex Sans', sans-serif;
        font-weight: 500;
        color: #8899a6;
        background-color: transparent;
        border-radius: 8px;
        padding: 0.75rem 1.5rem;
    }
    
    .stTabs [aria-selected="true"] {
        background: linear-gradient(90deg, #ff9500, #ff7b00) !important;
        color: white !important;
    }
    
    .sidebar .stRadio > label {
        font-size: 0.95rem;
    }
    
    div[data-testid="stExpander"] {
        background: #1a2634;
        border: 1px solid #2d3e50;
        border-radius: 12px;
    }
    
    .stProgress > div > div {
        background: linear-gradient(90deg, #ff9500, #00d4aa);
    }
    
    .upload-zone {
        border: 2px dashed #2d3e50;
        border-radius: 16px;
        padding: 3rem;
        text-align: center;
        background: rgba(26,38,52,0.5);
        transition: all 0.3s ease;
    }
    
    .upload-zone:hover {
        border-color: #ff9500;
        background: rgba(255,149,0,0.05);
    }
    
    .ai-response {
        background: linear-gradient(145deg, #0f1419 0%, #1a2634 100%);
        border: 1px solid #2d3e50;
        border-radius: 16px;
        padding: 1.5rem;
        margin-top: 1rem;
    }
    
    .ai-response h4 {
        color: #00d4aa;
        font-family: 'IBM Plex Sans', sans-serif;
        margin-bottom: 1rem;
    }
    
    hr {
        border-color: #2d3e50;
        margin: 2rem 0;
    }
    
    .stDownloadButton > button {
        background: linear-gradient(90deg, #00d4aa, #00b894) !important;
        font-family: 'IBM Plex Sans', sans-serif;
        font-weight: 600;
    }
    
    .stDownloadButton > button:hover {
        background: linear-gradient(90deg, #00b894, #00a67d) !important;
    }
</style>
""", unsafe_allow_html=True)

# Initialize session state
if 'assessment_data' not in st.session_state:
    st.session_state.assessment_data = {}
if 'ct_responses' not in st.session_state:
    st.session_state.ct_responses = {}
if 'ga_responses' not in st.session_state:
    st.session_state.ga_responses = {}
if 'ai_analysis' not in st.session_state:
    st.session_state.ai_analysis = None
if 'document_content' not in st.session_state:
    st.session_state.document_content = None

# Assessment Questionnaire Data
CONTROL_TOWER_DOMAINS = {
    "Landing Zone Architecture": {
        "weight": 0.20,
        "questions": [
            {
                "id": "ct_lz_1",
                "question": "What is the current state of your AWS multi-account strategy?",
                "options": {
                    "No multi-account strategy exists": 1,
                    "Basic account separation (dev/prod)": 2,
                    "Organizational units defined but not enforced": 3,
                    "Well-defined OU structure with some automation": 4,
                    "Mature multi-account with full automation": 5
                }
            },
            {
                "id": "ct_lz_2",
                "question": "How are new AWS accounts currently provisioned?",
                "options": {
                    "Manual creation through console": 1,
                    "Semi-automated with some scripts": 2,
                    "Automated but no standardization": 3,
                    "Account Factory with basic customization": 4,
                    "Fully automated Account Factory with AFT": 5
                }
            },
            {
                "id": "ct_lz_3",
                "question": "What is your current baseline configuration management approach?",
                "options": {
                    "No baseline configurations": 1,
                    "Manual baseline application": 2,
                    "Partial IaC implementation": 3,
                    "Full IaC with some drift detection": 4,
                    "Complete IaC with automated drift remediation": 5
                }
            },
            {
                "id": "ct_lz_4",
                "question": "How many AWS accounts are currently in scope for Control Tower migration?",
                "options": {
                    "1-10 accounts": 5,
                    "11-50 accounts": 4,
                    "51-100 accounts": 3,
                    "101-300 accounts": 2,
                    "300+ accounts": 1
                }
            }
        ]
    },
    "Governance & Guardrails": {
        "weight": 0.25,
        "questions": [
            {
                "id": "ct_gov_1",
                "question": "What is your current approach to preventive controls?",
                "options": {
                    "No preventive controls in place": 1,
                    "Basic IAM policies only": 2,
                    "SCPs defined but not comprehensive": 3,
                    "Comprehensive SCPs with some automation": 4,
                    "Full SCP hierarchy with automated enforcement": 5
                }
            },
            {
                "id": "ct_gov_2",
                "question": "How do you currently handle detective controls?",
                "options": {
                    "No detective controls": 1,
                    "Basic CloudTrail logging": 2,
                    "CloudTrail + some Config rules": 3,
                    "Comprehensive Config rules with alerts": 4,
                    "Full detective controls with auto-remediation": 5
                }
            },
            {
                "id": "ct_gov_3",
                "question": "What is your compliance posture for regulatory requirements?",
                "options": {
                    "No compliance framework": 1,
                    "Awareness of requirements only": 2,
                    "Partial compliance implementation": 3,
                    "Most compliance requirements met": 4,
                    "Full compliance with continuous monitoring": 5
                }
            },
            {
                "id": "ct_gov_4",
                "question": "How are tagging policies enforced across accounts?",
                "options": {
                    "No tagging strategy": 1,
                    "Tagging guidelines but no enforcement": 2,
                    "Some automated tagging checks": 3,
                    "Tag policies with partial enforcement": 4,
                    "Mandatory tagging with automated enforcement": 5
                }
            }
        ]
    },
    "Security & Identity": {
        "weight": 0.20,
        "questions": [
            {
                "id": "ct_sec_1",
                "question": "What is your current identity management approach for AWS?",
                "options": {
                    "Individual IAM users per account": 1,
                    "Federated access with basic setup": 2,
                    "AWS SSO with basic configuration": 3,
                    "IAM Identity Center with permission sets": 4,
                    "Fully integrated IAM Identity Center with ABAC": 5
                }
            },
            {
                "id": "ct_sec_2",
                "question": "How is cross-account access currently managed?",
                "options": {
                    "Shared credentials": 1,
                    "IAM users with cross-account roles": 2,
                    "Role assumption with basic trust policies": 3,
                    "Well-defined role hierarchy": 4,
                    "Automated role management with least privilege": 5
                }
            },
            {
                "id": "ct_sec_3",
                "question": "What is your security baseline for new accounts?",
                "options": {
                    "No security baseline": 1,
                    "Manual security configuration": 2,
                    "Partial security automation": 3,
                    "Security baseline with some gaps": 4,
                    "Comprehensive automated security baseline": 5
                }
            }
        ]
    },
    "Networking & Connectivity": {
        "weight": 0.15,
        "questions": [
            {
                "id": "ct_net_1",
                "question": "What is your current network architecture approach?",
                "options": {
                    "Individual VPCs with no connectivity": 1,
                    "VPC peering for some accounts": 2,
                    "Transit Gateway with basic setup": 3,
                    "Hub-spoke with Network Firewall": 4,
                    "Full network architecture with inspection": 5
                }
            },
            {
                "id": "ct_net_2",
                "question": "How is centralized DNS and resolution handled?",
                "options": {
                    "Individual Route53 zones per account": 1,
                    "Some shared hosted zones": 2,
                    "Route53 Resolver with basic rules": 3,
                    "Centralized DNS with some automation": 4,
                    "Fully automated centralized DNS management": 5
                }
            },
            {
                "id": "ct_net_3",
                "question": "What is your approach to network connectivity to on-premises?",
                "options": {
                    "Individual VPN connections per account": 1,
                    "Shared VPN with manual routing": 2,
                    "Direct Connect with basic setup": 3,
                    "Direct Connect Gateway with transit": 4,
                    "Full hybrid connectivity with redundancy": 5
                }
            }
        ]
    },
    "Operations & Monitoring": {
        "weight": 0.20,
        "questions": [
            {
                "id": "ct_ops_1",
                "question": "How is centralized logging currently implemented?",
                "options": {
                    "No centralized logging": 1,
                    "CloudTrail to individual S3 buckets": 2,
                    "Centralized logging bucket exists": 3,
                    "Log aggregation with some analysis": 4,
                    "Full log aggregation with SIEM integration": 5
                }
            },
            {
                "id": "ct_ops_2",
                "question": "What is your approach to cost management across accounts?",
                "options": {
                    "Individual account billing review": 1,
                    "Consolidated billing only": 2,
                    "Cost allocation tags defined": 3,
                    "Cost management with budgets and alerts": 4,
                    "Full FinOps with optimization automation": 5
                }
            },
            {
                "id": "ct_ops_3",
                "question": "How mature is your operational runbook and automation?",
                "options": {
                    "No documented runbooks": 1,
                    "Basic runbooks exist": 2,
                    "Runbooks with some automation": 3,
                    "Comprehensive runbooks with SSM": 4,
                    "Full automation with self-healing": 5
                }
            }
        ]
    }
}

GOLDEN_ARCHITECTURE_DOMAINS = {
    "Compute & Runtime": {
        "weight": 0.25,
        "questions": [
            {
                "id": "ga_comp_1",
                "question": "What is your current serverless compute adoption level?",
                "options": {
                    "No serverless adoption": 1,
                    "Experimental Lambda usage": 2,
                    "Some production Lambda workloads": 3,
                    "Significant serverless footprint": 4,
                    "Serverless-first strategy implemented": 5
                }
            },
            {
                "id": "ga_comp_2",
                "question": "How are Lambda functions organized and deployed?",
                "options": {
                    "Manual deployments through console": 1,
                    "CLI-based deployments": 2,
                    "SAM or Serverless Framework basics": 3,
                    "Full CI/CD with SAM/CDK": 4,
                    "GitOps with automated testing and deployment": 5
                }
            },
            {
                "id": "ga_comp_3",
                "question": "What is your approach to Lambda layers and shared code?",
                "options": {
                    "No shared code strategy": 1,
                    "Copy-paste code sharing": 2,
                    "Some Lambda layers in use": 3,
                    "Layer versioning with dependencies": 4,
                    "Automated layer management with CI/CD": 5
                }
            },
            {
                "id": "ga_comp_4",
                "question": "How do you handle container-based serverless (Fargate)?",
                "options": {
                    "No Fargate usage": 1,
                    "Experimental Fargate deployments": 2,
                    "Some Fargate in production": 3,
                    "Fargate with ECS/EKS integration": 4,
                    "Full serverless container strategy": 5
                }
            }
        ]
    },
    "API & Integration": {
        "weight": 0.20,
        "questions": [
            {
                "id": "ga_api_1",
                "question": "What is your API Gateway adoption level?",
                "options": {
                    "No API Gateway usage": 1,
                    "Basic REST APIs": 2,
                    "REST APIs with authorization": 3,
                    "HTTP APIs with advanced features": 4,
                    "Full API management with versioning": 5
                }
            },
            {
                "id": "ga_api_2",
                "question": "How are event-driven architectures implemented?",
                "options": {
                    "No event-driven patterns": 1,
                    "Basic SNS/SQS usage": 2,
                    "EventBridge for some integrations": 3,
                    "Event-driven with patterns defined": 4,
                    "Full event mesh with EventBridge": 5
                }
            },
            {
                "id": "ga_api_3",
                "question": "What is your approach to Step Functions and orchestration?",
                "options": {
                    "No workflow orchestration": 1,
                    "Basic Step Functions experimentation": 2,
                    "Step Functions for some workflows": 3,
                    "Standard and Express workflows": 4,
                    "Full orchestration with error handling": 5
                }
            }
        ]
    },
    "Data & Storage": {
        "weight": 0.20,
        "questions": [
            {
                "id": "ga_data_1",
                "question": "What serverless database services are you using?",
                "options": {
                    "Traditional RDS only": 1,
                    "Some DynamoDB experimentation": 2,
                    "DynamoDB in production": 3,
                    "DynamoDB with Aurora Serverless": 4,
                    "Full serverless data tier strategy": 5
                }
            },
            {
                "id": "ga_data_2",
                "question": "How do you handle data lake and analytics serverlessly?",
                "options": {
                    "No serverless analytics": 1,
                    "Basic S3 data storage": 2,
                    "Athena for ad-hoc queries": 3,
                    "Data lake with Glue and Athena": 4,
                    "Full serverless analytics platform": 5
                }
            },
            {
                "id": "ga_data_3",
                "question": "What is your caching strategy for serverless?",
                "options": {
                    "No caching strategy": 1,
                    "Basic ElastiCache usage": 2,
                    "DAX for DynamoDB": 3,
                    "Multi-layer caching approach": 4,
                    "Full caching with CloudFront + DAX": 5
                }
            }
        ]
    },
    "Security & Compliance": {
        "weight": 0.20,
        "questions": [
            {
                "id": "ga_sec_1",
                "question": "How are serverless function permissions managed?",
                "options": {
                    "Overly permissive IAM roles": 1,
                    "Basic role separation": 2,
                    "Least privilege attempted": 3,
                    "IAM roles with resource policies": 4,
                    "Fine-grained IAM with automated review": 5
                }
            },
            {
                "id": "ga_sec_2",
                "question": "What is your approach to secrets management?",
                "options": {
                    "Hardcoded secrets": 1,
                    "Environment variables only": 2,
                    "Parameter Store usage": 3,
                    "Secrets Manager with rotation": 4,
                    "Full secrets management with audit": 5
                }
            },
            {
                "id": "ga_sec_3",
                "question": "How do you handle API security and authentication?",
                "options": {
                    "No API authentication": 1,
                    "API keys only": 2,
                    "Cognito with basic setup": 3,
                    "Cognito with custom authorizers": 4,
                    "Full identity with WAF protection": 5
                }
            }
        ]
    },
    "Observability & DevOps": {
        "weight": 0.15,
        "questions": [
            {
                "id": "ga_obs_1",
                "question": "What is your serverless monitoring approach?",
                "options": {
                    "Basic CloudWatch metrics only": 1,
                    "CloudWatch with some custom metrics": 2,
                    "CloudWatch Insights enabled": 3,
                    "X-Ray tracing implemented": 4,
                    "Full observability with custom dashboards": 5
                }
            },
            {
                "id": "ga_obs_2",
                "question": "How mature is your serverless CI/CD pipeline?",
                "options": {
                    "Manual deployments": 1,
                    "Basic CI/CD pipeline": 2,
                    "Multi-stage deployments": 3,
                    "Canary deployments implemented": 4,
                    "Full GitOps with automated rollback": 5
                }
            },
            {
                "id": "ga_obs_3",
                "question": "What is your approach to serverless testing?",
                "options": {
                    "No automated testing": 1,
                    "Unit tests only": 2,
                    "Unit + integration tests": 3,
                    "Full testing with local emulation": 4,
                    "Comprehensive testing with chaos engineering": 5
                }
            }
        ]
    }
}


def calculate_domain_score(responses: Dict[str, int], domain_questions: List[Dict]) -> float:
    """Calculate weighted score for a domain."""
    if not responses:
        return 0.0
    
    total_score = 0
    answered = 0
    
    for question in domain_questions:
        if question["id"] in responses:
            total_score += responses[question["id"]]
            answered += 1
    
    if answered == 0:
        return 0.0
    
    return (total_score / (answered * 5)) * 100


def calculate_overall_score(responses: Dict[str, int], domains: Dict) -> Dict[str, float]:
    """Calculate overall assessment score."""
    domain_scores = {}
    weighted_total = 0
    
    for domain_name, domain_data in domains.items():
        score = calculate_domain_score(responses, domain_data["questions"])
        domain_scores[domain_name] = score
        weighted_total += score * domain_data["weight"]
    
    domain_scores["Overall"] = weighted_total
    return domain_scores


def get_maturity_level(score: float) -> tuple:
    """Get maturity level based on score."""
    if score >= 80:
        return "Optimized", "score-high"
    elif score >= 60:
        return "Managed", "score-medium"
    elif score >= 40:
        return "Developing", "score-medium"
    elif score >= 20:
        return "Initial", "score-low"
    else:
        return "Ad-hoc", "score-low"


def render_score_gauge(score: float, label: str):
    """Render a visual score gauge."""
    level, css_class = get_maturity_level(score)
    
    color = "#00d4aa" if score >= 60 else "#ff9500" if score >= 40 else "#ff6b6b"
    
    st.markdown(f"""
    <div class="metric-card">
        <div class="metric-value" style="color: {color};">{score:.0f}%</div>
        <div class="metric-label">{label}</div>
        <div class="score-badge {css_class}" style="margin-top: 0.5rem;">{level}</div>
    </div>
    """, unsafe_allow_html=True)


def call_claude_api(prompt: str, system_prompt: str = None) -> str:
    """Call Claude API for AI-driven analysis."""
    try:
        import anthropic
        
        api_key = os.environ.get("ANTHROPIC_API_KEY")
        if not api_key:
            return "⚠️ ANTHROPIC_API_KEY not configured. Please set the environment variable to enable AI analysis."
        
        client = anthropic.Anthropic(api_key=api_key)
        
        messages = [{"role": "user", "content": prompt}]
        
        response = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=4096,
            system=system_prompt or "You are an expert AWS Solutions Architect specializing in Control Tower migrations and serverless golden architectures. Provide detailed, actionable insights.",
            messages=messages
        )
        
        return response.content[0].text
        
    except ImportError:
        return "⚠️ Anthropic library not installed. Run: pip install anthropic"
    except Exception as e:
        return f"⚠️ AI Analysis Error: {str(e)}"


def extract_document_content(uploaded_file) -> str:
    """Extract content from uploaded documents."""
    content = ""
    file_type = uploaded_file.type
    
    try:
        if file_type == "application/pdf":
            try:
                import pypdf
                pdf_reader = pypdf.PdfReader(io.BytesIO(uploaded_file.read()))
                for page in pdf_reader.pages:
                    content += page.extract_text() + "\n"
            except ImportError:
                try:
                    import PyPDF2
                    pdf_reader = PyPDF2.PdfReader(io.BytesIO(uploaded_file.read()))
                    for page in pdf_reader.pages:
                        content += page.extract_text() + "\n"
                except ImportError:
                    content = "PDF extraction requires pypdf or PyPDF2 library."
                    
        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            try:
                import docx
                doc = docx.Document(io.BytesIO(uploaded_file.read()))
                for para in doc.paragraphs:
                    content += para.text + "\n"
            except ImportError:
                content = "DOCX extraction requires python-docx library."
                
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(uploaded_file.read()))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n"
            except ImportError:
                content = "PPTX extraction requires python-pptx library."
                
        elif file_type in ["text/plain", "application/json"]:
            content = uploaded_file.read().decode("utf-8")
            
    except Exception as e:
        content = f"Error extracting content: {str(e)}"
    
    return content


def generate_assessment_report(ct_scores: Dict, ga_scores: Dict, ai_analysis: str = None) -> str:
    """Generate comprehensive assessment report."""
    report = f"""
# AWS Enterprise Assessment Report
## Generated: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}

---

# Executive Summary

This assessment evaluates your organization's readiness for AWS Control Tower migration and Golden Architecture (Serverless) adoption.

---

# Control Tower Migration Assessment

## Overall Maturity Score: {ct_scores.get('Overall', 0):.1f}%
### Maturity Level: {get_maturity_level(ct_scores.get('Overall', 0))[0]}

### Domain Scores

"""
    
    for domain, score in ct_scores.items():
        if domain != "Overall":
            level, _ = get_maturity_level(score)
            report += f"- **{domain}**: {score:.1f}% ({level})\n"
    
    report += f"""

---

# Golden Architecture (Serverless) Assessment

## Overall Maturity Score: {ga_scores.get('Overall', 0):.1f}%
### Maturity Level: {get_maturity_level(ga_scores.get('Overall', 0))[0]}

### Domain Scores

"""
    
    for domain, score in ga_scores.items():
        if domain != "Overall":
            level, _ = get_maturity_level(score)
            report += f"- **{domain}**: {score:.1f}% ({level})\n"
    
    if ai_analysis:
        report += f"""

---

# AI-Driven Analysis & Recommendations

{ai_analysis}
"""
    
    report += f"""

---

# Next Steps

1. Review identified gaps and prioritize remediation
2. Develop detailed implementation roadmap
3. Establish governance framework
4. Plan phased migration approach
5. Define success metrics and KPIs

---

*Report generated by AWS Enterprise Assessment Platform*
"""
    
    return report


# Main Application
def main():
    # Header
    st.markdown("""
    <div class="main-header">
        <h1>🏗️ AWS Enterprise Assessment Platform</h1>
        <p>AI-Driven Control Tower Migration & Golden Architecture Assessment</p>
    </div>
    """, unsafe_allow_html=True)
    
    # Sidebar
    with st.sidebar:
        st.markdown("### 📋 Assessment Configuration")
        
        assessment_mode = st.radio(
            "Select Assessment Mode",
            ["🎯 Questionnaire-Based", "📄 Document Analysis", "🔄 Hybrid Assessment"],
            index=0
        )
        
        st.markdown("---")
        
        st.markdown("### 🏢 Organization Details")
        org_name = st.text_input("Organization Name", placeholder="Enter organization name")
        assessor_name = st.text_input("Assessor Name", placeholder="Enter your name")
        
        st.markdown("---")
        
        st.markdown("### ⚙️ AI Configuration")
        enable_ai = st.checkbox("Enable AI Analysis", value=True)
        
        if enable_ai:
            ai_detail_level = st.select_slider(
                "Analysis Detail Level",
                options=["Brief", "Standard", "Comprehensive"],
                value="Standard"
            )
        
        st.markdown("---")
        
        st.markdown("### 📊 Quick Stats")
        ct_complete = len(st.session_state.ct_responses)
        ct_total = sum(len(d["questions"]) for d in CONTROL_TOWER_DOMAINS.values())
        ga_complete = len(st.session_state.ga_responses)
        ga_total = sum(len(d["questions"]) for d in GOLDEN_ARCHITECTURE_DOMAINS.values())
        
        st.metric("Control Tower Progress", f"{ct_complete}/{ct_total}")
        st.metric("Golden Arch Progress", f"{ga_complete}/{ga_total}")
    
    # Main content tabs
    tab1, tab2, tab3, tab4, tab5 = st.tabs([
        "🎛️ Control Tower Assessment",
        "⚡ Golden Architecture Assessment", 
        "📄 Document Analysis",
        "🤖 AI Insights",
        "📊 Reports & Export"
    ])
    
    # Tab 1: Control Tower Assessment
    with tab1:
        st.markdown("## Control Tower Migration Readiness Assessment")
        st.markdown("Evaluate your organization's readiness for AWS Control Tower adoption.")
        
        for domain_name, domain_data in CONTROL_TOWER_DOMAINS.items():
            with st.expander(f"📁 {domain_name}", expanded=False):
                st.markdown(f'<div class="domain-header">{domain_name} (Weight: {domain_data["weight"]*100:.0f}%)</div>', unsafe_allow_html=True)
                
                for question in domain_data["questions"]:
                    q_id = question["id"]
                    st.markdown(f"**{question['question']}**")
                    
                    options = list(question["options"].keys())
                    
                    current_idx = 0
                    if q_id in st.session_state.ct_responses:
                        current_val = st.session_state.ct_responses[q_id]
                        for idx, (opt, val) in enumerate(question["options"].items()):
                            if val == current_val:
                                current_idx = idx
                                break
                    
                    selected = st.radio(
                        f"Select response for {q_id}",
                        options,
                        index=current_idx,
                        key=f"radio_{q_id}",
                        label_visibility="collapsed"
                    )
                    
                    st.session_state.ct_responses[q_id] = question["options"][selected]
                    st.markdown("---")
        
        # Calculate and display scores
        if st.session_state.ct_responses:
            st.markdown("### 📈 Current Assessment Scores")
            ct_scores = calculate_overall_score(st.session_state.ct_responses, CONTROL_TOWER_DOMAINS)
            
            cols = st.columns(3)
            with cols[1]:
                render_score_gauge(ct_scores["Overall"], "Overall Readiness")
            
            st.markdown("### Domain Breakdown")
            domain_cols = st.columns(len(ct_scores) - 1)
            for idx, (domain, score) in enumerate(ct_scores.items()):
                if domain != "Overall":
                    with domain_cols[idx % len(domain_cols)]:
                        render_score_gauge(score, domain[:20])
    
    # Tab 2: Golden Architecture Assessment
    with tab2:
        st.markdown("## Golden Architecture (Serverless) Assessment")
        st.markdown("Evaluate your serverless architecture maturity and readiness.")
        
        for domain_name, domain_data in GOLDEN_ARCHITECTURE_DOMAINS.items():
            with st.expander(f"📁 {domain_name}", expanded=False):
                st.markdown(f'<div class="domain-header">{domain_name} (Weight: {domain_data["weight"]*100:.0f}%)</div>', unsafe_allow_html=True)
                
                for question in domain_data["questions"]:
                    q_id = question["id"]
                    st.markdown(f"**{question['question']}**")
                    
                    options = list(question["options"].keys())
                    
                    current_idx = 0
                    if q_id in st.session_state.ga_responses:
                        current_val = st.session_state.ga_responses[q_id]
                        for idx, (opt, val) in enumerate(question["options"].items()):
                            if val == current_val:
                                current_idx = idx
                                break
                    
                    selected = st.radio(
                        f"Select response for {q_id}",
                        options,
                        index=current_idx,
                        key=f"radio_{q_id}",
                        label_visibility="collapsed"
                    )
                    
                    st.session_state.ga_responses[q_id] = question["options"][selected]
                    st.markdown("---")
        
        # Calculate and display scores
        if st.session_state.ga_responses:
            st.markdown("### 📈 Current Assessment Scores")
            ga_scores = calculate_overall_score(st.session_state.ga_responses, GOLDEN_ARCHITECTURE_DOMAINS)
            
            cols = st.columns(3)
            with cols[1]:
                render_score_gauge(ga_scores["Overall"], "Overall Maturity")
            
            st.markdown("### Domain Breakdown")
            domain_cols = st.columns(len(ga_scores) - 1)
            for idx, (domain, score) in enumerate(ga_scores.items()):
                if domain != "Overall":
                    with domain_cols[idx % len(domain_cols)]:
                        render_score_gauge(score, domain[:20])
    
    # Tab 3: Document Analysis
    with tab3:
        st.markdown("## 📄 Document-Based Assessment")
        st.markdown("Upload existing documentation for AI-powered analysis and gap identification.")
        
        st.markdown("""
        <div class="upload-zone">
            <h3>📤 Upload Assessment Documents</h3>
            <p>Supported formats: PDF, DOCX, PPTX, TXT, JSON</p>
        </div>
        """, unsafe_allow_html=True)
        
        uploaded_files = st.file_uploader(
            "Upload your documents",
            type=["pdf", "docx", "pptx", "txt", "json"],
            accept_multiple_files=True,
            label_visibility="collapsed"
        )
        
        if uploaded_files:
            st.markdown("### 📑 Uploaded Documents")
            
            all_content = ""
            for file in uploaded_files:
                with st.expander(f"📄 {file.name}"):
                    content = extract_document_content(file)
                    all_content += f"\n\n--- Document: {file.name} ---\n{content}"
                    st.text_area("Content Preview", content[:2000] + "..." if len(content) > 2000 else content, height=200)
            
            st.session_state.document_content = all_content
            
            if st.button("🤖 Analyze Documents with AI", type="primary"):
                with st.spinner("Analyzing documents..."):
                    analysis_prompt = f"""
Analyze the following AWS-related documentation and provide a comprehensive assessment for:

1. **Control Tower Migration Readiness**:
   - Current multi-account strategy
   - Governance and guardrails maturity
   - Security and identity posture
   - Network architecture readiness
   - Operational capabilities

2. **Golden Architecture (Serverless) Readiness**:
   - Serverless compute adoption
   - API and integration patterns
   - Data and storage architecture
   - Security implementation
   - Observability and DevOps maturity

For each area, provide:
- Current state assessment (1-5 scale)
- Key gaps identified
- Specific recommendations
- Priority ranking (High/Medium/Low)

Documentation Content:
{all_content[:15000]}

Provide structured, actionable insights for enterprise implementation.
"""
                    
                    analysis = call_claude_api(analysis_prompt)
                    st.session_state.ai_analysis = analysis
                    
                    st.markdown("### 🔍 AI Analysis Results")
                    st.markdown(f"""
                    <div class="ai-response">
                        {analysis}
                    </div>
                    """, unsafe_allow_html=True)
    
    # Tab 4: AI Insights
    with tab4:
        st.markdown("## 🤖 AI-Driven Insights & Recommendations")
        
        ct_scores = calculate_overall_score(st.session_state.ct_responses, CONTROL_TOWER_DOMAINS) if st.session_state.ct_responses else {}
        ga_scores = calculate_overall_score(st.session_state.ga_responses, GOLDEN_ARCHITECTURE_DOMAINS) if st.session_state.ga_responses else {}
        
        if ct_scores or ga_scores:
            st.markdown("### 📊 Assessment Summary for AI Analysis")
            
            col1, col2 = st.columns(2)
            with col1:
                st.markdown("**Control Tower Scores:**")
                for domain, score in ct_scores.items():
                    st.write(f"- {domain}: {score:.1f}%")
            with col2:
                st.markdown("**Golden Architecture Scores:**")
                for domain, score in ga_scores.items():
                    st.write(f"- {domain}: {score:.1f}%")
        
        analysis_type = st.selectbox(
            "Select Analysis Type",
            [
                "🎯 Gap Analysis & Prioritization",
                "🗺️ Implementation Roadmap",
                "⚠️ Risk Assessment",
                "💰 Cost-Benefit Analysis",
                "🏗️ Architecture Recommendations",
                "📋 Compliance Mapping"
            ]
        )
        
        custom_context = st.text_area(
            "Additional Context (optional)",
            placeholder="Enter any additional context about your organization, constraints, or specific requirements...",
            height=100
        )
        
        if st.button("🚀 Generate AI Analysis", type="primary"):
            with st.spinner("Generating comprehensive analysis..."):
                
                scores_summary = f"""
Control Tower Assessment Scores:
{json.dumps(ct_scores, indent=2) if ct_scores else "Not completed"}

Golden Architecture Assessment Scores:
{json.dumps(ga_scores, indent=2) if ga_scores else "Not completed"}
"""
                
                analysis_prompts = {
                    "🎯 Gap Analysis & Prioritization": f"""
Based on the following assessment scores, provide a detailed gap analysis:

{scores_summary}

Additional Context: {custom_context or "None provided"}

Please provide:
1. **Critical Gaps** (Score < 40%): Immediate attention required
2. **Moderate Gaps** (Score 40-60%): Short-term focus areas
3. **Minor Gaps** (Score 60-80%): Optimization opportunities
4. **Strengths** (Score > 80%): Areas to leverage

For each gap:
- Root cause analysis
- Business impact assessment
- Recommended remediation steps
- Estimated effort (T-shirt sizing)
- Dependencies and prerequisites
""",
                    "🗺️ Implementation Roadmap": f"""
Based on the assessment scores, create a detailed implementation roadmap:

{scores_summary}

Additional Context: {custom_context or "None provided"}

Provide a phased approach:

**Phase 1: Foundation (Months 1-3)**
- Quick wins and prerequisites
- Resource requirements
- Key milestones

**Phase 2: Core Implementation (Months 4-6)**
- Control Tower deployment steps
- Serverless architecture migration
- Governance framework

**Phase 3: Optimization (Months 7-12)**
- Advanced features
- Automation enhancement
- Continuous improvement

Include dependencies, risks, and success metrics for each phase.
""",
                    "⚠️ Risk Assessment": f"""
Perform a comprehensive risk assessment based on:

{scores_summary}

Additional Context: {custom_context or "None provided"}

Analyze risks across:
1. **Technical Risks**: Architecture, integration, complexity
2. **Operational Risks**: Skills gap, change management, BAU impact
3. **Security Risks**: Compliance, data protection, access control
4. **Financial Risks**: Cost overruns, resource constraints
5. **Timeline Risks**: Dependencies, scope creep

For each risk provide:
- Probability (High/Medium/Low)
- Impact (High/Medium/Low)
- Mitigation strategies
- Contingency plans
""",
                    "💰 Cost-Benefit Analysis": f"""
Provide a cost-benefit analysis for the migration:

{scores_summary}

Additional Context: {custom_context or "None provided"}

Include:
1. **Cost Categories**:
   - Implementation costs (labor, tooling, training)
   - Migration costs (data transfer, parallel running)
   - Ongoing operational costs
   - Opportunity costs

2. **Benefits**:
   - Operational efficiency gains
   - Security posture improvement
   - Compliance automation
   - Cost optimization potential
   - Innovation enablement

3. **ROI Calculation Framework**
4. **Break-even Analysis**
5. **TCO Comparison**
""",
                    "🏗️ Architecture Recommendations": f"""
Provide detailed architecture recommendations based on:

{scores_summary}

Additional Context: {custom_context or "None provided"}

Cover:
1. **Control Tower Architecture**:
   - OU structure recommendation
   - Account factory design
   - Guardrail selection
   - Network topology

2. **Golden Architecture (Serverless)**:
   - Reference architecture patterns
   - Service selection guidance
   - Integration patterns
   - Data architecture

3. **Cross-cutting Concerns**:
   - Security architecture
   - Observability strategy
   - CI/CD pipeline design
   - Disaster recovery

Include architecture diagrams descriptions and implementation notes.
""",
                    "📋 Compliance Mapping": f"""
Map compliance requirements to implementation:

{scores_summary}

Additional Context: {custom_context or "None provided"}

Provide mapping for common frameworks:
1. **SOC 2**: Control mapping to CT guardrails
2. **PCI DSS**: Relevant controls and implementation
3. **HIPAA**: Security and privacy requirements
4. **GDPR**: Data protection considerations
5. **AWS Well-Architected**: Pillar alignment

For each:
- Current compliance gaps
- Control Tower controls that address requirements
- Additional controls needed
- Evidence collection automation
"""
                }
                
                selected_prompt = analysis_prompts.get(analysis_type, analysis_prompts["🎯 Gap Analysis & Prioritization"])
                
                analysis_result = call_claude_api(selected_prompt)
                st.session_state.ai_analysis = analysis_result
                
                st.markdown("### 📋 Analysis Results")
                st.markdown(f"""
                <div class="ai-response">
                    <h4>🤖 AI-Generated {analysis_type}</h4>
                </div>
                """, unsafe_allow_html=True)
                st.markdown(analysis_result)
    
    # Tab 5: Reports & Export
    with tab5:
        st.markdown("## 📊 Assessment Reports & Export")
        
        ct_scores = calculate_overall_score(st.session_state.ct_responses, CONTROL_TOWER_DOMAINS) if st.session_state.ct_responses else {}
        ga_scores = calculate_overall_score(st.session_state.ga_responses, GOLDEN_ARCHITECTURE_DOMAINS) if st.session_state.ga_responses else {}
        
        # Summary Dashboard
        st.markdown("### 📈 Assessment Dashboard")
        
        col1, col2, col3, col4 = st.columns(4)
        
        with col1:
            ct_overall = ct_scores.get("Overall", 0)
            st.markdown(f"""
            <div class="metric-card">
                <div class="metric-value" style="color: {'#00d4aa' if ct_overall >= 60 else '#ff9500' if ct_overall >= 40 else '#ff6b6b'};">
                    {ct_overall:.0f}%
                </div>
                <div class="metric-label">Control Tower</div>
            </div>
            """, unsafe_allow_html=True)
        
        with col2:
            ga_overall = ga_scores.get("Overall", 0)
            st.markdown(f"""
            <div class="metric-card">
                <div class="metric-value" style="color: {'#00d4aa' if ga_overall >= 60 else '#ff9500' if ga_overall >= 40 else '#ff6b6b'};">
                    {ga_overall:.0f}%
                </div>
                <div class="metric-label">Golden Architecture</div>
            </div>
            """, unsafe_allow_html=True)
        
        with col3:
            combined = (ct_overall + ga_overall) / 2 if ct_overall and ga_overall else ct_overall or ga_overall
            st.markdown(f"""
            <div class="metric-card">
                <div class="metric-value" style="color: {'#00d4aa' if combined >= 60 else '#ff9500' if combined >= 40 else '#ff6b6b'};">
                    {combined:.0f}%
                </div>
                <div class="metric-label">Combined Score</div>
            </div>
            """, unsafe_allow_html=True)
        
        with col4:
            level, _ = get_maturity_level(combined)
            st.markdown(f"""
            <div class="metric-card">
                <div class="metric-value" style="font-size: 1.5rem; color: #ff9500;">
                    {level}
                </div>
                <div class="metric-label">Maturity Level</div>
            </div>
            """, unsafe_allow_html=True)
        
        st.markdown("---")
        
        # Export Options
        st.markdown("### 📥 Export Options")
        
        col1, col2, col3 = st.columns(3)
        
        with col1:
            if st.button("📄 Generate Full Report", type="primary"):
                report = generate_assessment_report(ct_scores, ga_scores, st.session_state.ai_analysis)
                st.session_state.generated_report = report
                st.success("Report generated successfully!")
        
        with col2:
            if st.button("📊 Export Raw Data (JSON)"):
                export_data = {
                    "metadata": {
                        "generated_at": datetime.now().isoformat(),
                        "organization": org_name,
                        "assessor": assessor_name
                    },
                    "control_tower": {
                        "responses": st.session_state.ct_responses,
                        "scores": ct_scores
                    },
                    "golden_architecture": {
                        "responses": st.session_state.ga_responses,
                        "scores": ga_scores
                    },
                    "ai_analysis": st.session_state.ai_analysis
                }
                st.download_button(
                    "⬇️ Download JSON",
                    json.dumps(export_data, indent=2),
                    file_name=f"aws_assessment_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json",
                    mime="application/json"
                )
        
        with col3:
            if 'generated_report' in st.session_state:
                st.download_button(
                    "⬇️ Download Report (MD)",
                    st.session_state.generated_report,
                    file_name=f"aws_assessment_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.md",
                    mime="text/markdown"
                )
        
        # Display generated report
        if 'generated_report' in st.session_state:
            st.markdown("---")
            st.markdown("### 📋 Generated Report Preview")
            with st.expander("View Full Report", expanded=True):
                st.markdown(st.session_state.generated_report)
        
        # Detailed Domain Analysis
        st.markdown("---")
        st.markdown("### 📊 Detailed Domain Analysis")
        
        tab_ct, tab_ga = st.tabs(["Control Tower Domains", "Golden Architecture Domains"])
        
        with tab_ct:
            if ct_scores:
                for domain, score in ct_scores.items():
                    if domain != "Overall":
                        level, css_class = get_maturity_level(score)
                        with st.expander(f"{domain}: {score:.1f}% ({level})"):
                            st.progress(score / 100)
                            
                            # Show question responses
                            domain_data = CONTROL_TOWER_DOMAINS.get(domain, {})
                            for q in domain_data.get("questions", []):
                                if q["id"] in st.session_state.ct_responses:
                                    response_val = st.session_state.ct_responses[q["id"]]
                                    response_text = [k for k, v in q["options"].items() if v == response_val][0]
                                    st.markdown(f"**Q:** {q['question']}")
                                    st.markdown(f"**A:** {response_text} (Score: {response_val}/5)")
                                    st.markdown("---")
        
        with tab_ga:
            if ga_scores:
                for domain, score in ga_scores.items():
                    if domain != "Overall":
                        level, css_class = get_maturity_level(score)
                        with st.expander(f"{domain}: {score:.1f}% ({level})"):
                            st.progress(score / 100)
                            
                            # Show question responses
                            domain_data = GOLDEN_ARCHITECTURE_DOMAINS.get(domain, {})
                            for q in domain_data.get("questions", []):
                                if q["id"] in st.session_state.ga_responses:
                                    response_val = st.session_state.ga_responses[q["id"]]
                                    response_text = [k for k, v in q["options"].items() if v == response_val][0]
                                    st.markdown(f"**Q:** {q['question']}")
                                    st.markdown(f"**A:** {response_text} (Score: {response_val}/5)")
                                    st.markdown("---")


if __name__ == "__main__":
    main()
